from pathlib import Path


def parse_pdf(path: Path) -> str:
    import pypdf
    reader = pypdf.PdfReader(str(path))
    return "\n".join(
        page.extract_text() or "" for page in reader.pages
    )


def parse_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"## Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(slides)


def parse_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def parse_doc(path: Path) -> str:
    """旧版 .doc：调用 LibreOffice 转成文本"""
    import subprocess, tempfile
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "txt:Text", "--outdir", tmp, str(path)],
            check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        txt_path = Path(tmp) / (path.stem + ".txt")
        return txt_path.read_text(encoding="utf-8", errors="ignore")


def parse_document(path: Path) -> str:
    ext = path.suffix.lower()
    print(f"[parser] 解析文档: {path.name}")
    if ext == ".pdf":
        return parse_pdf(path)
    if ext == ".pptx":
        return parse_pptx(path)
    if ext == ".docx":
        return parse_docx(path)
    if ext in {".ppt", ".doc"}:
        return parse_doc(path)
    raise ValueError(f"不支持的文件类型: {ext}")
